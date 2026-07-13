from __future__ import annotations

import json
import os
import re
import sys
import urllib.error
import urllib.request
from pathlib import Path

import typer
from rich.console import Console

app = typer.Typer(add_completion=False)
console = Console()


DB_DIR = Path(os.environ.get("KIS_DB", ".kis_db"))
EMBED_MODEL = os.environ.get("KIS_EMBED", "BAAI/bge-small-en-v1.5")
LLM_MODEL = os.environ.get("KIS_MODEL", "qwen2.5:1.5b")
OLLAMA_URL = os.environ.get("KIS_OLLAMA", "http://localhost:11434")
CHUNK_CHARS = int(os.environ.get("KIS_CHUNK", "1200"))
CHUNK_OVERLAP = int(os.environ.get("KIS_OVERLAP", "150"))
TOP_K = int(os.environ.get("KIS_TOPK", "4"))
MIN_SIMILARITY = float(os.environ.get("KIS_MINSIM", "0.40"))

TEXT_EXT = {".txt", ".md", ".markdown", ".rst", ".log", ".csv",
            ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".go",
            ".rs", ".rb", ".sh", ".json", ".yaml", ".yml", ".html", ".css"}


def is_prose_block(text: str) -> bool:
    words = text.split()
    lc = sum(1 for w in words if w.isalpha() and w.islower() and len(w) >= 3)
    return lc >= 3 or bool(re.match(r"\s*(Figure|Table)\s+\d+", text))


def parse_pdf(path: Path):
    import fitz
    doc = fitz.open(path)
    out = []
    for i, page in enumerate(doc):
        text = "\n".join(b[4] for b in page.get_text("blocks")
                         if is_prose_block(b[4])).strip()
        if text:
            out.append((text, f"p.{i + 1}"))
    doc.close()
    return out


def parse_docx(path: Path):
    import docx
    text = "\n".join(p.text for p in docx.Document(path).paragraphs if p.text.strip())
    return [(text, "")] if text.strip() else []


def parse_pptx(path: Path):
    from pptx import Presentation
    out = []
    for i, slide in enumerate(Presentation(path).slides):
        parts = [sh.text for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
        if parts:
            out.append(("\n".join(parts), f"slide {i + 1}"))
    return out


def parse_xlsx(path: Path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        rows = ["\t".join("" if c is None else str(c) for c in row)
                for row in ws.iter_rows(values_only=True)]
        text = "\n".join(r for r in rows if r.strip())
        if text.strip():
            out.append((text, f"sheet: {ws.title}"))
    wb.close()
    return out


def parse_email(path: Path):
    import email
    from email import policy
    msg = email.message_from_bytes(path.read_bytes(), policy=policy.default)
    body = msg.get_body(preferencelist=("plain", "html"))
    text = body.get_content() if body else ""
    header = f"From: {msg['from']}\nTo: {msg['to']}\nDate: {msg['date']}\nSubject: {msg['subject']}"
    subj = (msg["subject"] or "email")[:60]
    return [(header + "\n\n" + text, f"email: {subj}")] if text.strip() else []


def parse_mbox(path: Path):
    import mailbox
    out = []
    for msg in mailbox.mbox(str(path)):
        payload = msg.get_payload(decode=True)
        text = payload.decode(errors="replace") if payload else str(msg.get_payload())
        subj = (msg["subject"] or "email")[:60]
        if text.strip():
            out.append((f"Subject: {msg['subject']}\n\n{text}", f"email: {subj}"))
    return out


def parse_text(path: Path):
    text = path.read_text(encoding="utf-8", errors="replace")
    return [(text, "")] if text.strip() else []


def parse_file(path: Path):
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return parse_pdf(path)
        if ext == ".docx":
            return parse_docx(path)
        if ext == ".pptx":
            return parse_pptx(path)
        if ext == ".xlsx":
            return parse_xlsx(path)
        if ext == ".eml":
            return parse_email(path)
        if ext == ".mbox":
            return parse_mbox(path)
        if ext in TEXT_EXT:
            return parse_text(path)
    except Exception as e:
        console.print(f"[yellow]skip[/] {path.name}: {e}")
    return None


SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".eml", ".mbox"} | TEXT_EXT


def chunk_text(text: str):
    text = text.strip()
    if len(text) <= CHUNK_CHARS:
        return [text] if text else []
    chunks, start = [], 0
    while start < len(text):
        end = start + CHUNK_CHARS
        if end < len(text):
            brk = text.rfind(" ", start + CHUNK_CHARS - CHUNK_OVERLAP, end)
            if brk != -1:
                end = brk
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = end - CHUNK_OVERLAP
    return [c for c in chunks if c]


def build_records(path: Path, root: Path):
    segments = parse_file(path)
    if not segments:
        return []
    source = str(path.relative_to(root))
    mtime = path.stat().st_mtime_ns
    records = []
    for text, loc in segments:
        for chunk in chunk_text(text):
            records.append({"source": source, "loc": loc, "mtime": mtime,
                            "chunk_index": len(records), "text": chunk})
    return records


_embedder = None


def embedder():
    global _embedder
    if _embedder is None:
        from fastembed import TextEmbedding
        console.print(f"[dim]loading embedder {EMBED_MODEL} (first run downloads ~130MB)…[/]")
        _embedder = TextEmbedding(model_name=EMBED_MODEL)
    return _embedder


def embed_passages(texts):
    return [v.tolist() for v in embedder().embed(texts)]


def embed_query(text):
    return next(embedder().query_embed(text)).tolist()


def open_table(create_from=None):
    import lancedb
    db = lancedb.connect(DB_DIR)
    if "chunks" in db.list_tables().tables:
        return db.open_table("chunks")
    if create_from:
        return db.create_table("chunks", data=create_from)
    return None


def index_stats():
    table = open_table()
    if table is None:
        return None
    sources = table.to_arrow()["source"].to_pylist()
    counts = {}
    for s in sources:
        counts[s] = counts.get(s, 0) + 1
    return len(sources), counts


def reset_index():
    import lancedb
    db = lancedb.connect(DB_DIR)
    if "chunks" not in db.list_tables().tables:
        return False
    db.drop_table("chunks")
    return True


def find_missing(folder):
    root = Path(folder).expanduser().resolve()
    stats = index_stats()
    if stats is None:
        return []
    _, counts = stats
    return [(s, n) for s, n in counts.items() if not (root / s).exists()]


def delete_sources(sources):
    table = open_table()
    if table is None or not sources:
        return 0
    for s in sources:
        table.delete(f"source = '{s.replace(chr(39), chr(39) * 2)}'")
    if table.count_rows():
        from lancedb.index import FTS
        table.create_index("text", config=FTS(), replace=True)
    return len(sources)


def ollama_status():
    try:
        req = urllib.request.Request(f"{OLLAMA_URL}/api/tags")
        with urllib.request.urlopen(req, timeout=5) as resp:
            models = [m["name"] for m in json.loads(resp.read()).get("models", [])]
    except Exception:
        return False, f"Ollama not reachable at {OLLAMA_URL} — start it (see ollama.com)."
    want = LLM_MODEL.split(":")[0]
    if any(m == LLM_MODEL or m.split(":")[0] == want for m in models):
        return True, f"Ollama running · model '{LLM_MODEL}' ready"
    return False, f"Ollama running, but model '{LLM_MODEL}' not pulled — run: ollama pull {LLM_MODEL}"


def find_ingestable(folder: str):
    root = Path(folder).expanduser().resolve()
    if not root.is_dir():
        raise ValueError(f"Not a folder: {root}")
    files = [p for p in root.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED]
    if not files:
        raise ValueError(f"No supported files under {root}")
    return root, files


def ingest_files(files, root, progress=None):
    table = open_table()
    indexed = {}
    if table is not None:
        t = table.to_arrow()
        if "mtime" in t.column_names:
            indexed = dict(zip(t["source"].to_pylist(), t["mtime"].to_pylist()))
    added = updated = skipped = 0
    for n, path in enumerate(files, 1):
        source = str(path.relative_to(root))
        if indexed.get(source) == path.stat().st_mtime_ns:
            skipped += 1
            if progress:
                progress(n, len(files), source, 0)
            continue
        records = build_records(path, root)
        if not records:
            if progress:
                progress(n, len(files), source, 0)
            continue
        for r, vec in zip(records, embed_passages([r["text"] for r in records])):
            r["vector"] = vec
        if source in indexed and table is not None:
            table.delete(f"source = '{source.replace(chr(39), chr(39) * 2)}'")
            updated += 1
        if table is None:
            table = open_table(create_from=records)
        else:
            table.add(records)
        added += len(records)
        if progress:
            progress(n, len(files), source, len(records))
    if added:
        from lancedb.index import FTS
        table.create_index("text", config=FTS(), replace=True)
    return {"added": added, "updated": updated, "skipped": skipped,
            "empty": table is None}


@app.command()
def ingest(folder: str):
    try:
        root, files = find_ingestable(folder)
    except ValueError as e:
        console.print(f"[red]{e}[/]")
        raise typer.Exit(1)
    console.print(f"Found [bold]{len(files)}[/] files. Parsing + embedding…")

    def show(n, total, source, n_chunks):
        if n_chunks:
            console.print(f"  [dim][{n}/{total}][/] [green]+{n_chunks}[/] {source}")

    stats = ingest_files(files, root, progress=show)
    if stats["empty"]:
        console.print("[yellow]Nothing could be extracted from those files.[/]")
        raise typer.Exit(1)
    console.print(f"[bold green]Done.[/] +{stats['added']} chunks "
                  f"({stats['updated']} files updated, {stats['skipped']} unchanged) "
                  f"in {DB_DIR}/")


@app.command()
def status():
    stats = index_stats()
    if stats is None:
        console.print("[yellow]No index yet. Run `ingest` first.[/]")
        raise typer.Exit(1)
    total, counts = stats
    console.print(f"[bold]{total}[/] chunks from [bold]{len(counts)}[/] files.")
    for src in sorted(counts):
        console.print(f"  {src}  [dim]({counts[src]} chunks)[/]")


@app.command()
def reset(yes: bool = typer.Option(False, "--yes", "-y",
                                   help="Skip the confirmation prompt.")):
    stats = index_stats()
    if stats is None:
        console.print("[yellow]No index to reset.[/]")
        raise typer.Exit(0)
    total, counts = stats
    if not yes and not typer.confirm(
            f"Delete the entire index ({total} chunks from {len(counts)} files) "
            f"in {DB_DIR}?"):
        console.print("Aborted.")
        raise typer.Exit(0)
    reset_index()
    console.print(f"[bold green]Index reset.[/] Removed {total} chunks from {DB_DIR}/")


@app.command()
def prune(folder: str,
          yes: bool = typer.Option(False, "--yes", "-y",
                                   help="Skip the confirmation prompt.")):
    if index_stats() is None:
        console.print("[yellow]No index yet. Run `ingest` first.[/]")
        raise typer.Exit(1)
    missing = find_missing(folder)
    if not missing:
        console.print("[green]Nothing to prune[/] — every indexed file still exists.")
        raise typer.Exit(0)
    for src, n in missing:
        console.print(f"  [red]-{n}[/] {src}")
    chunks = sum(n for _, n in missing)
    if not yes and not typer.confirm(
            f"Remove {chunks} chunks from {len(missing)} missing file(s)?"):
        console.print("Aborted.")
        raise typer.Exit(0)
    delete_sources([src for src, _ in missing])
    console.print(f"[bold green]Pruned.[/] Removed {chunks} chunks "
                  f"from {len(missing)} file(s).")


def retrieve(query: str, k: int = TOP_K):
    table = open_table()
    if table is None:
        return None, []
    hits = (table.search(embed_query(query)).metric("cosine").limit(k).to_list())
    for h in hits:
        h["similarity"] = 1 - h["_distance"]
    seen = {h["text"] for h in hits}
    try:
        for h in table.search(query, query_type="fts").limit(k).to_list():
            if h["text"] not in seen:
                h["similarity"] = None
                hits.append(h)
                seen.add(h["text"])
    except Exception:
        pass
    return table, hits


def ollama_chat(system: str, user: str) -> str:
    payload = {"model": LLM_MODEL, "stream": False,
               "options": {"temperature": 0},
               "messages": [{"role": "system", "content": system},
                            {"role": "user", "content": user}]}
    req = urllib.request.Request(f"{OLLAMA_URL}/api/chat",
                                 data=json.dumps(payload).encode(),
                                 headers={"Content-Type": "application/json"})
    try:
        with urllib.request.urlopen(req, timeout=300) as resp:
            return json.loads(resp.read())["message"]["content"].strip()
    except urllib.error.URLError as e:
        if not hasattr(e, "code"):
            console.print(f"[red]Can't reach Ollama at {OLLAMA_URL}.[/] "
                          "Is it installed and running? See: https://ollama.com")
            raise typer.Exit(1)
        console.print(f"[red]Ollama error {e.code}:[/] {e.read().decode(errors='replace')}")
        console.print(f"If the model is missing, run:  [bold]ollama pull {LLM_MODEL}[/]")
        raise typer.Exit(1)


SYSTEM_PROMPT = (
    "You answer questions using ONLY the numbered context passages given to you. "
    "Read the passages and give a short, direct answer based only on facts that "
    "appear in them — never use outside knowledge. "
    "If the question is ambiguous or could mean several different things, ask ONE "
    "short clarifying question instead of guessing. "
    "If none of the passages contain the answer, reply exactly: "
    "I couldn't find this in your documents."
)


def build_user_prompt(context: str, question: str) -> str:
    return (f"Context:\n{context}\n\n"
            f"Question: {question}\n\n"
            "Answer the question in one or two sentences using only the context "
            "above.\nAnswer:")


def format_context(hits) -> str:
    return "\n\n".join(
        f"[{i + 1}] {h['text']}\n(from {h['source']}{', ' + h['loc'] if h['loc'] else ''})"
        for i, h in enumerate(hits))


def print_sources(hits):
    console.print("[bold]Sources:[/]")
    for i, h in enumerate(hits):
        loc = f", {h['loc']}" if h["loc"] else ""
        score = f"~{h['similarity']:.2f}" if h["similarity"] is not None else "keyword"
        console.print(f"  [{i + 1}] {h['source']}{loc}  [dim]({score})[/]")


@app.command()
def ask(question: str):
    table, hits = retrieve(question)
    if table is None:
        console.print("[yellow]No index yet. Run `ingest` first.[/]")
        raise typer.Exit(1)

    if not hits or hits[0]["similarity"] < MIN_SIMILARITY:
        console.print("[yellow]I couldn't find this in your documents.[/]")
        if hits:
            console.print(f"[dim](closest match ~{hits[0]['similarity']:.2f}, "
                          f"below threshold {MIN_SIMILARITY})[/]")
        raise typer.Exit(0)

    answer = ollama_chat(SYSTEM_PROMPT, build_user_prompt(format_context(hits), question))
    console.print(f"\n{answer}\n", markup=False)
    print_sources(hits)


@app.command()
def where(term: str):
    table = open_table()
    if table is None:
        console.print("[yellow]No index yet. Run `ingest` first.[/]")
        raise typer.Exit(1)
    try:
        hits = table.search(term, query_type="fts").limit(20).to_list()
    except Exception:
        console.print("[yellow]No keyword index. Re-run `ingest` to build it.[/]")
        raise typer.Exit(1)
    if not hits:
        console.print(f"[yellow]'{term}' not found in your documents.[/]")
        raise typer.Exit(0)
    console.print(f"[bold]{len(hits)}[/] passage(s) mention '{term}':")
    for h in hits:
        loc = f", {h['loc']}" if h["loc"] else ""
        snippet = " ".join(h["text"].split())[:120]
        console.print(f"  [cyan]{h['source']}{loc}[/]: [dim]{snippet}…[/]")


INSIGHTS_PROMPT = (
    "You summarize what the numbered context passages say about the user's topic. "
    "Use ONLY the passages — never outside knowledge. Write a short synthesis "
    "(2-5 sentences) and cite the passages you used by their number, like [1]. "
    "If the passages don't actually discuss the topic, reply exactly: "
    "I couldn't find this in your documents."
)


@app.command()
def about(topic: str):
    table, hits = retrieve(topic, k=5)
    if table is None:
        console.print("[yellow]No index yet. Run `ingest` first.[/]")
        raise typer.Exit(1)
    if not hits or hits[0]["similarity"] < MIN_SIMILARITY:
        console.print(f"[yellow]I couldn't find '{topic}' in your documents.[/]")
        raise typer.Exit(0)
    hits = hits[:6]
    summary = ollama_chat(INSIGHTS_PROMPT,
                          f"Topic: {topic}\n\nContext:\n{format_context(hits)}\n\n"
                          "Summarize what the context says about the topic, "
                          "citing passages like [1].\nSummary:")
    console.print(f"\n{summary}\n", markup=False)
    docs = list(dict.fromkeys(h["source"] for h in hits))
    console.print(f"[bold]Across {len(docs)} document(s):[/] {', '.join(docs)}")
    print_sources(hits)


@app.command()
def selftest():
    assert chunk_text("hello world") == ["hello world"]
    assert chunk_text("   ") == []
    text = " ".join(f"word{i}" for i in range(2000))
    chunks = chunk_text(text)
    assert len(chunks) >= 2
    assert all(len(c) <= CHUNK_CHARS for c in chunks)
    assert chunks[0].split()[0] == "word0"
    assert chunks[-1].split()[-1] == "word1999"
    assert chunks[0].split()[-1] in chunks[1]
    assert is_prose_block("The model was trained on the development set.")
    assert is_prose_block("Figure 4: Illustrations of Fine-tuning BERT.")
    assert not is_prose_block("E[CLS] E1 E[SEP] EN T1 TN [CLS] Tok 1 Tok M")
    assert not is_prose_block("6 512 2048 8 64 64 0.1 0.1 100K 4.92 25.8 65")
    console.print("[green]selftest passed[/] — chunker + prose filter OK")


if __name__ == "__main__":
    app()
